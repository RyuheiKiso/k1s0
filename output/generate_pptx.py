#!/usr/bin/env python3
"""Generate PowerPoint presentation for developer profile analysis."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK = RGBColor(0x2C, 0x3E, 0x50)
BLUE = RGBColor(0x34, 0x98, 0xDB)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFB)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_para(tf, text, size=18, color=BLACK, bold=False, space_before=Pt(6), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.alignment = align
    return p


# ===== SLIDE 1: Cover =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK)
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "開発者心理分析レポート", 44, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
             "木曽 龍平（Ryuhei Kiso）— k1s0 プロジェクト創設者", 24, RGBColor(0xBD, 0xC3, 0xC7), False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.8), Inches(11), Inches(1.5),
             "対象：32歳 / エンジニア歴8年\n分析手法：コードベース・設計文書・哲学的記述の多角的心理学分析\n作成日：2026年1月30日",
             16, GRAY, False, PP_ALIGN.CENTER)

# ===== SLIDE 2: kiso_memo.txt =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "分析対象：philosophy/kiso_memo.txt 全10箇条", 28, WHITE, True)

items = [
    "1. システムにおいて一番の不確定要素は流動的な人員配置による開発体制（人間）である。",
    "2. 末端の開発者に倫理を期待してはいけない。幹の人間が責任を負うべきである。",
    "3. 一個人の拡大解釈を認めてはいけない。",
    "4. 生成AIでチェックするのはあくまで成果物であり、バグの本質は作成の過程にある。",
    "5. 自由とカオスは紙一重。",
    "6. 末端の開発者は未来を考えてくれないものである。",
    "7. 下のスキルに合わせる必要のない空間を提供したい。",
    "8. 挑戦する空間と安定した空間を提供したい。",
    "9. 部外者はシステムの現場を理解する気はないものである。",
    "10. 他者に期待してはいけない。",
]
tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), items[0], 17, BLACK)
for item in items[1:]:
    add_para(tf, item, 17, BLACK, space_before=Pt(8))

# ===== SLIDE 3: 精神分析 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "II. 精神分析的観点 — 昇華（Sublimation）", 28, WHITE, True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), "防衛機制としてのk1s0", 22, DARK, True)
tf = add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
    "「他者に期待してはいけない」は、繰り返し期待し裏切られた経験の反復強迫の果てに到達した防衛的結論。", 16, BLACK)
add_para(tf, "", 10)
add_para(tf, "対象関係論（クライン）：「良い対象」への信頼が破壊され、他者を信頼対象から除外する妄想-分裂ポジション的世界認識に意識的に移行。", 16)
add_para(tf, "", 10)
add_para(tf, "しかし病理的退行ではなく、意図的な設計判断として外在化。内的な傷つきをk1s0という外的構造物に昇華している。", 16)
add_para(tf, "", 10)
add_para(tf, "→ フロイト的に最も成熟した防衛機制「昇華」の典型例", 17, BLUE, True)

add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5), LIGHT_BG)
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5), "昇華のメカニズム", 20, DARK, True)
tf = add_text_box(slide, Inches(7.3), Inches(2.4), Inches(5), Inches(3.8),
    "内的苦痛（他者への失望）", 16, RED, True)
add_para(tf, "　　↓", 16, GRAY)
add_para(tf, "防衛的結論（期待の放棄）", 16, RED, True)
add_para(tf, "　　↓", 16, GRAY)
add_para(tf, "外在化（k1s0の設計思想へ）", 16, BLUE, True)
add_para(tf, "　　↓", 16, GRAY)
add_para(tf, "昇華（社会的に有用な創造物）", 16, GREEN, True)

# ===== SLIDE 4: アドラー =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "III. アドラー心理学 — 劣等感と補償・共同体感覚", 28, WHITE, True)

add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.4), "劣等感の源泉", 20, DARK, True)
tf = add_text_box(slide, Inches(1), Inches(2.3), Inches(5), Inches(1.5),
    "自身の能力ではなく「環境の無秩序」に対する劣等感。理想の開発環境が他者によって実現されないという深い無力感。", 16)

add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4), "補償行動", 20, DARK, True)
tf = add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(1.5),
    "k1s0の構築 = 社会的に有用な補償行動。個人的な怒りや失望を、他者も使えるプラットフォームとして社会に還元。", 16)

add_shape_bg(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(2.8), RGBColor(0xFD, 0xF2, 0xF2))
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.4), "共同体感覚との緊張", 20, RED, True)
tf = add_text_box(slide, Inches(1), Inches(5.1), Inches(11), Inches(1.8),
    "「他者に期待しない」と言いながら、膨大なドキュメントを書き、他者が使うことを前提として設計している。", 16)
add_para(tf, "", 10)
add_para(tf, "→ 矛盾ではなく「期待しないが、機会は提供する」という独自の共同体感覚", 17, RED, True)

# ===== SLIDE 5: マズロー =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "IV. マズローの欲求階層 — 下位欲求を飛び越える使命感", 28, WHITE, True)

rows = [
    ("自己実現", "「挑戦する空間と安定した空間を提供したい」", "活性化", GREEN),
    ("承認欲求", "直接的な言及なし", "抑圧", RED),
    ("所属と愛", "「部外者は理解する気はない」「他者に期待するな」", "放棄", RED),
    ("安全欲求", "「間違ったものを作らせない」「自由とカオスは紙一重」", "未充足", RED),
]
y = 1.5
for label, desc, status, scolor in rows:
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(1.2), scolor)
    add_text_box(slide, Inches(0.9), Inches(y + 0.1), Inches(2.3), Inches(1), label, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.5), Inches(y + 0.1), Inches(6.5), Inches(1), desc, 15)
    add_text_box(slide, Inches(10.5), Inches(y + 0.1), Inches(2), Inches(1), status, 16, scolor, True, PP_ALIGN.CENTER)
    y += 1.35

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(1),
    "→ マズロー晩年の修正理論「欲求階層の逆転」に該当。強い使命感が下位欲求を飛び越える。", 17, DARK, True)

# ===== SLIDE 6: ユング =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "V. ユング分析心理学 — ペルソナ・シャドウ・元型", 28, WHITE, True)

add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.4), "ペルソナ（外的人格）", 20, BLUE, True)
add_text_box(slide, Inches(1), Inches(2.3), Inches(5), Inches(1.5),
    "極めて整然とした秩序。11のLintルール、Clean Architectureの厳格な依存方向、managed/protected分離。", 16)

add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(2.5), RGBColor(0xFD, 0xF2, 0xF2))
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4), "シャドウ（影）", 20, RED, True)
add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(1.5),
    "自身の中にある秩序を壊したいという衝動。6言語ソロ実装という限界突破的行動がその表出。内的カオスを外的秩序で制御。", 16)

add_shape_bg(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(1), Inches(4.5), Inches(5), Inches(0.4), "老賢者（Senex）", 20, DARK, True)
add_text_box(slide, Inches(1), Inches(5.1), Inches(5), Inches(1.5),
    "哲学的箴言を書き記す行為。32歳にして40代の達観を持つ。", 16)

add_shape_bg(slide, Inches(7), Inches(4.3), Inches(5.5), Inches(2.5), LIGHT_BG)
add_text_box(slide, Inches(7.2), Inches(4.5), Inches(5), Inches(0.4), "永遠の少年（Puer Aeternus）", 20, DARK, True)
add_text_box(slide, Inches(7.2), Inches(5.1), Inches(5), Inches(1.5),
    "ソロで壮大なプラットフォームを構築する——無限の可能性を信じ、一人で世界を変えようとする衝動。", 16)

# ===== SLIDE 7: CBT =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "VI. 認知行動療法的観点 — スキーマ分析", 28, WHITE, True)

schemas = [
    ("他者不信", "「他者に期待してはいけない」\n「末端に倫理を期待するな」", "断絶と拒絶"),
    ("警戒", "「拡大解釈を認めてはいけない」\n「自由とカオスは紙一重」", "過剰警戒"),
    ("自己犠牲", "「幹の人間が責任を負うべき」", "他者志向"),
    ("厳密な基準", "520行の懸念事項\n24個のエッジケース文書化", "過剰警戒"),
]
y = 1.5
for name, evidence, category in schemas:
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(1.1), BLUE)
    add_text_box(slide, Inches(0.9), Inches(y + 0.15), Inches(2.3), Inches(0.8), name, 17, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.5), Inches(y + 0.05), Inches(6), Inches(1), evidence, 14)
    add_text_box(slide, Inches(10), Inches(y + 0.15), Inches(2.5), Inches(0.8), category, 15, GRAY, False, PP_ALIGN.CENTER)
    y += 1.2

add_shape_bg(slide, Inches(0.8), Inches(6), Inches(11.7), Inches(1.2), RGBColor(0xFD, 0xF2, 0xF2))
add_text_box(slide, Inches(1), Inches(6.1), Inches(11.3), Inches(1),
    "重要：これらは幼少期トラウマではなく「職業的信念体系」（8年の帰納的結論）である可能性が高い。\n不適応的スキーマを最も機能する領域（システム設計）に正確に配置 → 高い自己認識の表れ。", 16, RED)

# ===== SLIDE 8: 発達心理学 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "VII. 発達心理学 — 早すぎる生殖性と30歳の過渡期", 28, WHITE, True)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), "エリクソン発達段階", 20, DARK, True)
stages = [
    ("24-26歳", "新卒〜若手。理想の形成期"),
    ("26-28歳", "中堅。理想と現実の乖離を体感"),
    ("28-30歳", "リーダー的役割。「末端は未来を考えない」の実体験"),
    ("30-32歳", "基盤側へ。怒りから構築への転換"),
]
y = 2.2
for age, desc in stages:
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(1.8), Inches(0.7), BLUE)
    add_text_box(slide, Inches(0.9), Inches(y + 0.1), Inches(1.6), Inches(0.5), age, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.8), Inches(y + 0.1), Inches(3.5), Inches(0.5), desc, 14)
    y += 0.85

add_shape_bg(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.2), LIGHT_BG)
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.4), "核心的発見", 20, DARK, True)
tf = add_text_box(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(4),
    "32歳 = エリクソン第6段階（親密性 vs 孤立）の時期", 16, BLACK, True)
add_para(tf, "", 10)
add_para(tf, "しかし第6段階を飛び越え、第7段階「生殖性」（次世代への貢献）に突入している。", 16)
add_para(tf, "", 10)
add_para(tf, "レヴィンソンの「30歳の過渡期」：20代の人生構造への根本的問い直し。kiso_memo.txtはその爆発の記録。", 16)
add_para(tf, "", 10)
add_para(tf, "→ 「信頼しないが、与える」", 18, RED, True)
add_para(tf, "不信を前提としたまま生殖性を実現する独自の発達的統合。", 16)

# ===== SLIDE 9: ギフテッド =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "VIII. 認知的成熟度とギフテッド特性", 28, WHITE, True)

traits = [
    ("過度激動", "6言語を一人で実装するエネルギー量"),
    ("完璧主義", "520行の懸念事項、24のエッジケース"),
    ("非同期発達", "知的成熟は40代水準、社会的信頼は未発達"),
    ("存在論的孤独", "「他者に期待するな」——理解者の不在"),
    ("強い正義感", "「幹の人間が責任を負うべき」"),
    ("内的駆動力", "承認なしに大規模プロジェクトを推進"),
]
y = 1.5
for trait, evidence in traits:
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(2.2), Inches(0.8), BLUE)
    add_text_box(slide, Inches(0.9), Inches(y + 0.1), Inches(2), Inches(0.6), trait, 15, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.2), Inches(y + 0.1), Inches(4), Inches(0.6), evidence, 15)
    y += 0.9

add_shape_bg(slide, Inches(7.5), Inches(1.5), Inches(5), Inches(5.2), LIGHT_BG)
add_text_box(slide, Inches(7.7), Inches(1.7), Inches(4.6), Inches(0.4), "ポスト形式的思考", 20, DARK, True)
tf = add_text_box(slide, Inches(7.7), Inches(2.3), Inches(4.6), Inches(1.5),
    "「バグの本質は作成の過程にある」= メタ認知の高さ\n「自由とカオスは紙一重」= 弁証法的思考\n→ 通常30代後半以降に発達する能力", 15)

add_text_box(slide, Inches(7.7), Inches(4), Inches(4.6), Inches(0.4), "ドンブロフスキの積極的分離", 18, DARK, True)
tf = add_text_box(slide, Inches(7.7), Inches(4.6), Inches(4.6), Inches(2),
    "既存の価値体系を一度解体し、自分自身の価値体系を再構築する過程。\n\nkiso_memo.txt = 積極的分離の過程で生まれた再構築された価値体系の記録。", 15)

# ===== SLIDE 10: リスクと可能性 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8), "IX. リスクと可能性", 28, WHITE, True)

risks = [
    ("燃え尽き", "8年の失望＋ソロ開発負荷。30-35歳はバーンアウト頻発年齢帯"),
    ("親密性の回避固定化", "エリクソン第6段階を未解決のまま通過 → 孤立の慢性化"),
    ("過剰統制", "秩序への希求が自身の柔軟性を喪失させる"),
    ("委譲能力の未発達", "「自分でやった方が早い」が協働の可能性を閉じる"),
]
y = 1.5
for name, desc in risks:
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(5.5), Inches(0.9), RGBColor(0xFD, 0xF2, 0xF2))
    add_text_box(slide, Inches(0.9), Inches(y + 0.05), Inches(1.8), Inches(0.8), name, 14, RED, True)
    add_text_box(slide, Inches(2.8), Inches(y + 0.05), Inches(3.3), Inches(0.8), desc, 13)
    y += 1.0

opportunities = [
    ("安定した生産期", "30歳の過渡期を創造物で通過 → 33歳以降に安定"),
    ("メンターへの発展", "「提供したい」欲求が師弟関係に発展しうる"),
    ("思想の普遍化", "10箇条は組織論として普遍的。知的生産はまだ序盤"),
    ("信念の進化", "「期待するな」→「期待の仕方を変える」への進化余地"),
]
y = 1.5
for name, desc in opportunities:
    add_shape_bg(slide, Inches(7), Inches(y), Inches(5.5), Inches(0.9), RGBColor(0xEA, 0xFA, 0xEA))
    add_text_box(slide, Inches(7.1), Inches(y + 0.05), Inches(1.8), Inches(0.8), name, 14, GREEN, True)
    add_text_box(slide, Inches(9), Inches(y + 0.05), Inches(3.3), Inches(0.8), desc, 13)
    y += 1.0

# ===== SLIDE 11: 総合所見 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK)
add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8), "X. 総合所見", 32, WHITE, True, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(2.2), RGBColor(0x34, 0x49, 0x5E))
tf = add_text_box(slide, Inches(1.3), Inches(2), Inches(10.7), Inches(2),
    "この人物は、同年代より著しく高い認知的成熟を持ち、8年間の職業経験で蓄積した「人間の不完全さ」への洞察を、32歳にして既に体系的な設計哲学に昇華させている。その代償として対人的親密性の発達が抑制されているが、これは不可逆ではなく、k1s0の成功と協働者の出現によって開かれる余地がある。",
    17, WHITE)

add_text_box(slide, Inches(1), Inches(4.3), Inches(11.333), Inches(0.5), "心理的肖像（一文要約）", 20, GRAY, True, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(1), Inches(5), Inches(11.333), Inches(1.8), RGBColor(0x34, 0x49, 0x5E))
tf = add_text_box(slide, Inches(1.3), Inches(5.2), Inches(10.7), Inches(1.5),
    "繰り返し傷ついた信頼を、構造的秩序の創造という昇華によって回復不能な形で再構築しようとしている、高機能な孤独の建築家。\n\n人間を信頼しない。だからこそ、人間を助ける仕組みを作る。",
    20, WHITE, True, PP_ALIGN.CENTER)

# Save
prs.save("/home/user/k1s0/output/developer_profile_analysis.pptx")
print("PPTX saved.")
